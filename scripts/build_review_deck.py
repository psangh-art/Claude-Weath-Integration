#!/usr/bin/env python3
"""Build a PowerPoint review deck from the latest pipeline run: sorted by layout
then chart, ONE SLIDE PER INVESTMENT, showing the captured chart image alongside
everything the pipeline knows about it (live price, master-sheet holdings/alert
levels, channel read, TradingView alerts) — with loud red/amber flags wherever a
chart or alert is MISSING, so the whole run can be eyeballed in PowerPoint Online.

Data sources (all produced by run_full_pipeline.js):
  scripts/layout_manifest_tmp.json    charts + screenshots + live prices
  scripts/alerts_manifest_tmp.json    live TradingView alerts
  scripts/channel_results_tmp.json    OCR channel-boundary reads
  ~/Downloads/Stocks_Buy_Strategy.xlsx  master sheet ('Investments' tab)

Usage: python build_review_deck.py [out.pptx]
  (default out: ~/Downloads/Investment_Review_Deck.pptx)
"""
import html as _html
import json
import os
import sys
import urllib.parse
from datetime import datetime

import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw, ImageFont

from ticker_normalize import normalize, master_tickers_match

if sys.platform == "win32":
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
from config import CFG, downloads_dir, purge_old_versions
DOWNLOADS = downloads_dir()
MASTER_PATH = os.path.join(DOWNLOADS, CFG['masterWorkbook'])
# The Investment Production Centre front end serves these two: an in-Chrome
# gallery view of the deck and a machine-readable summary for the output bay.
PIPELINE_APP_DIR = os.path.join(SCRIPT_DIR, 'pipeline_app')
GALLERY_PATH = os.path.join(PIPELINE_APP_DIR, 'review_deck.html')
SUMMARY_PATH = os.path.join(PIPELINE_APP_DIR, 'review_deck_summary.json')
TV_LAYOUT_URL = CFG['tvLayoutUrlTemplate']

# Master 'Investments' columns (same map as update_master_sheet.py)
# +1 vs the pre-2026-07-16 layout: a 'Marked Up' column was inserted at Investments!B.
COL_SHARE_NAME, COL_TICKER, COL_HOLDINGS, COL_TARGET = 3, 4, 5, 7
COL_ALERT_LOW, COL_ALERT_LOW_SOURCE, COL_ALERT_HIGH = 13, 14, 16

NAVY = RGBColor(0x1F, 0x38, 0x64)
RED = RGBColor(0xCC, 0x00, 0x00)
AMBER = RGBColor(0xB9, 0x77, 0x0E)
GREEN = RGBColor(0x27, 0x62, 0x21)
GREY = RGBColor(0x60, 0x60, 0x60)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)


def load_json(name):
    path = os.path.join(SCRIPT_DIR, name)
    if not os.path.exists(path):
        return None
    with open(path, encoding='utf-8') as f:
        return json.load(f)


def load_master_index():
    """ticker(str, master form) -> dict of the master row's review-relevant cells."""
    if not os.path.exists(MASTER_PATH):
        return {}
    wb = openpyxl.load_workbook(MASTER_PATH, data_only=False)
    ws = wb['Investments']
    index = {}
    for r in range(1, ws.max_row + 1):
        t = ws.cell(row=r, column=COL_TICKER).value
        if not isinstance(t, str) or not t.strip() or t.strip().upper() == 'TICKER':
            continue  # skip blanks and the header row
        index[t.strip().upper()] = {
            'row': r,
            'share_name': ws.cell(row=r, column=COL_SHARE_NAME).value,
            'holdings': ws.cell(row=r, column=COL_HOLDINGS).value,
            'target': ws.cell(row=r, column=COL_TARGET).value,
            'alert_low': ws.cell(row=r, column=COL_ALERT_LOW).value,
            'alert_low_source': ws.cell(row=r, column=COL_ALERT_LOW_SOURCE).value,
            'alert_high': ws.cell(row=r, column=COL_ALERT_HIGH).value,
        }
    return index


def find_master(index, chart_ticker):
    norm = normalize(chart_ticker)
    if not norm or not norm['master_ticker']:
        return None, None
    for key, row in index.items():
        if master_tickers_match(key, norm['master_ticker']):
            return key, row
    return norm['master_ticker'], None


def alerts_for(alerts, chart_ticker):
    """TradingView alerts whose symbol's post-colon part matches this chart ticker."""
    out = []
    want = chart_ticker.strip().upper()
    for a in alerts or []:
        sym = (a.get('symbol') or '').split(':')[-1].strip().upper()
        if sym == want or sym.rstrip('.') == want.rstrip('.'):
            out.append(a)
    return out


def fmt_pounds(v):
    if isinstance(v, (int, float)):
        return f'£{v:,.0f}'
    return '—'


def fmt_num(v):
    if isinstance(v, (int, float)):
        return f'{v:,.2f}'.rstrip('0').rstrip('.')
    return '—'


def fmt_ts(v):
    """Human-readable capture timestamp: an ISO string like
    '2026-07-17T16:29:48.104Z' -> '17 Jul 2026, 16:29'. Leaves anything it
    can't parse untouched, and returns '—' for empties."""
    if not v:
        return '—'
    if not isinstance(v, str):
        return str(v)
    s = v.strip().replace('Z', '+00:00')
    try:
        dt = datetime.fromisoformat(s)
    except ValueError:
        # Fall back to just the date+minute portion of a raw ISO string.
        return v.replace('T', ' ')[:16]
    return dt.strftime('%d %b %Y, %H:%M')


def add_text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT):
    """lines: list of (text, size_pt, bold, colour) tuples, one paragraph each."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for text, size, bold, colour in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = colour
    return box


def add_flag_banner(slide, y, text):
    box = add_text(slide, Inches(8.55), y, Inches(4.5), Inches(0.35),
                   [(text, 13, True, RGBColor(0xFF, 0xFF, 0xFF))])
    box.fill.solid()
    box.fill.fore_color.rgb = RED
    return y + Inches(0.42)


def add_picture_fitted(slide, path, x, y, max_w, max_h):
    with Image.open(path) as im:
        pw, ph = im.size
    scale = min(max_w / pw, max_h / ph)
    w, h = int(pw * scale), int(ph * scale)
    slide.shapes.add_picture(path, x, y + Emu(int((max_h - h) / 2)), width=Emu(w), height=Emu(h))


# ── Alert-level overlay ────────────────────────────────────────────────────
# The user checks Alert Low/High by eye across the deck, so draw the detected
# levels straight onto each chart image: a horizontal line at the price's pixel
# row (from channel_detect's axis fit price = a*y + b) with a labelled tag. Low
# is green, High is orange — both distinct from the yellow trend lines / blue
# channel already on the chart. Annotated copies go to a throwaway dir; the pptx
# embeds the bytes at add_picture time so the files only need to exist briefly.
_ANNOTATED_DIR = os.path.join(SCRIPT_DIR, 'pipeline_app', '_annotated_charts')
_LEVEL_COLOURS = {'low': (0x00, 0xE6, 0x76), 'high': (0xFF, 0x91, 0x00)}


def _level_font(size):
    for name in ('arialbd.ttf', 'arial.ttf', 'DejaVuSans-Bold.ttf'):
        try:
            return ImageFont.truetype(name, size)
        except OSError:
            continue
    return ImageFont.load_default()


def annotate_chart_levels(src_path, channel):
    """Return a path to a copy of src_path with the detected Alert Low/High drawn
    on, or src_path unchanged when there's no axis fit / no level to draw."""
    if not channel:
        return src_path
    a, b = channel.get('axis_a'), channel.get('axis_b')
    if not a:  # axis couldn't be read this run — nothing to place a level against
        return src_path
    levels = []
    if isinstance(channel.get('lower'), (int, float)):
        levels.append(('low', 'Low', float(channel['lower'])))
    if isinstance(channel.get('upper'), (int, float)):
        levels.append(('high', 'High', float(channel['upper'])))
    if not levels:
        return src_path
    try:
        im = Image.open(src_path).convert('RGB')
    except Exception:
        return src_path
    w, h = im.size
    draw = ImageDraw.Draw(im)
    font = _level_font(max(13, int(h * 0.024)))
    lw = max(2, int(h * 0.004))
    sw = max(1, lw // 2)
    for key, label, price in levels:
        y = int(round((price - b) / a))
        if not (0 <= y < h):     # level sits off the visible frame — skip its line
            continue
        colour = _LEVEL_COLOURS[key]
        draw.line([(0, y), (w, y)], fill=colour, width=lw)
        tag = f'{label} {price:.2f}'
        l, t, r, bo = draw.textbbox((0, 0), tag, font=font)
        th = bo - t
        pad = max(3, int(th * 0.25))
        # No background box — draw coloured text with a thin dark stroke so it
        # stays legible over the chart. Sit it just above the line.
        ty = min(max(0, y - th - pad - lw), h - th - pad)
        draw.text((pad, ty - t), tag, fill=colour, font=font,
                  stroke_width=sw, stroke_fill=(0, 0, 0))
    os.makedirs(_ANNOTATED_DIR, exist_ok=True)
    out = os.path.join(_ANNOTATED_DIR, os.path.basename(src_path))
    try:
        im.save(out)
    except Exception:
        return src_path
    return out


VERDICT_COLOURS = {'Buy candidate': GREEN, 'Hold': NAVY, 'Watch': AMBER, 'Avoid': RED}


def chart_slide(prs, chart, layout_name, master_key, master, channel, tv_alerts, analyst=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    ticker = chart.get('ticker') or '?'
    desc = chart.get('description') or ''
    add_text(slide, Inches(0.35), Inches(0.15), Inches(9.5), Inches(0.55),
             [(f'{ticker} — {desc}', 24, True, NAVY)])
    add_text(slide, Inches(10.0), Inches(0.22), Inches(3.0), Inches(0.4),
             [(layout_name, 12, False, GREY)], align=PP_ALIGN.RIGHT)

    # Chart image (left, large) or a loud MISSING placeholder
    img = chart.get('screenshot')
    img_x, img_y = Inches(0.35), Inches(0.85)
    img_w, img_h = Inches(8.0), Inches(6.3)
    if img and os.path.exists(img):
        # Draw the detected Alert Low/High straight onto the chart so the user can
        # confirm the levels at a glance (green = low, orange = high).
        add_picture_fitted(slide, annotate_chart_levels(img, channel), img_x, img_y, img_w, img_h)
    else:
        box = add_text(slide, img_x, Inches(3.2), img_w, Inches(1.2),
                       [('⚠ MISSING CHART', 36, True, RGBColor(0xFF, 0xFF, 0xFF)),
                        (chart.get('error') or 'no screenshot captured this run', 14, False,
                         RGBColor(0xFF, 0xFF, 0xFF))],
                       align=PP_ALIGN.CENTER)
        box.fill.solid()
        box.fill.fore_color.rgb = RED

    # Info column (right)
    x, w = Inches(8.55), Inches(4.5)
    y = Inches(0.85)

    if master is None:
        y = add_flag_banner(slide, y, '⚠ NOT IN MASTER SHEET')
    has_tv_alert = bool(tv_alerts)
    has_master_alert = master is not None and isinstance(master.get('alert_low'), (int, float))
    if not has_tv_alert and not has_master_alert:
        y = add_flag_banner(slide, y, '⚠ NO ALERTS (TradingView or master)')

    lines = [(f"Live price: {fmt_num(chart.get('price'))}", 15, True, NAVY),
             (f"  captured {fmt_ts(chart.get('priceCheckedAt'))}", 10, False, GREY)]

    if master:
        lines += [
            ('Master sheet (Investments)', 13, True, NAVY),
            (f"  Share name: {master.get('share_name') or '—'}", 12, False, GREY),
            (f"  Holdings: {fmt_pounds(master.get('holdings'))}    "
             f"Target: {fmt_pounds(master.get('target'))}", 12, False, GREY),
            (f"  Alert Low: {fmt_num(master.get('alert_low'))} "
             f"({master.get('alert_low_source') or 'no source'})", 12,
             not has_master_alert, RED if not has_master_alert else GREEN),
            (f"  Alert High: {fmt_num(master.get('alert_high'))}", 12, False, GREY),
        ]

    lines.append(('Channel read (OCR)', 13, True, NAVY))
    if channel and channel.get('kind') not in (None, 'rejected'):
        lines.append((f"  {channel['kind']}: lower {fmt_num(channel.get('lower'))}, "
                      f"upper {fmt_num(channel.get('upper'))}", 12, False, GREEN))
    elif channel:
        lines.append((f"  rejected: {channel.get('reason') or 'no reliable read'}", 12, False, AMBER))
    else:
        lines.append(('  not attempted this run', 12, False, GREY))

    # Analyst view (written by the investment-analyst agent into
    # scripts/analyst_notes.json; rendered whenever present for this ticker)
    if analyst:
        verdict = analyst.get('verdict') or 'Watch'
        colour = VERDICT_COLOURS.get(verdict, NAVY)
        lines.append(('Analyst view', 13, True, NAVY))
        buy = analyst.get('buy_price')
        basis = analyst.get('buy_basis') or ''
        buy_txt = f" — buy {fmt_num(buy)}" + (f' ({basis})' if basis and buy is not None else '')
        lines.append((f'  {verdict}{buy_txt if buy is not None else ""}', 12, True, colour))
        if analyst.get('note'):
            lines.append((f"  {analyst['note']}", 10, False, GREY))

    lines.append((f'TradingView alerts ({len(tv_alerts)})', 13, True,
                  NAVY if has_tv_alert else RED))
    for a in tv_alerts[:6]:
        status = 'active' if a.get('active') else 'INACTIVE'
        lines.append((f"  {fmt_num(a.get('targetPrice'))} {a.get('conditionType') or ''} "
                      f"[{status}] exp {str(a.get('expiration') or '—')[:10]}", 11, False,
                      GREY if a.get('active') else AMBER))
    if len(tv_alerts) > 6:
        lines.append((f'  … and {len(tv_alerts) - 6} more', 11, False, GREY))

    add_text(slide, x, y, w, Inches(7.3) - y, lines)


def section_slide(prs, title, subtitle=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = add_text(slide, Inches(0.8), Inches(3.0), Inches(11.7), Inches(1.5),
                   [(title, 40, True, RGBColor(0xFF, 0xFF, 0xFF))] +
                   ([(subtitle, 16, False, RGBColor(0xD0, 0xD8, 0xE8))] if subtitle else []),
                   align=PP_ALIGN.CENTER)
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY


def _asset_url(abs_path):
    """The front-end server proxies whitelisted images at /asset?p=<abs path>."""
    return 'asset?p=' + urllib.parse.quote(abs_path or '', safe='')


def _e(v):
    return _html.escape('' if v is None else str(v))


def write_gallery(path, charts, alerts, channels, master_index, stats):
    """Write an in-Chrome HTML view of the review deck — one card per chart,
    grouped by layout, with the same flags the .pptx carries. Shares the
    Investment Production Centre's dark-slate/amber design language."""
    def card(chart):
        ticker = chart.get('ticker') or '?'
        _, master = find_master(master_index, ticker)
        channel = channels.get(ticker)
        tv_alerts = alerts_for(alerts, ticker)
        img = chart.get('screenshot')
        cid = chart.get('chartId')
        has_master_alert = master is not None and isinstance(master.get('alert_low'), (int, float))
        flags = []
        if not (img and os.path.exists(img)):
            flags.append('<span class="flag bad">Missing chart</span>')
        if not tv_alerts and not has_master_alert:
            flags.append('<span class="flag bad">No alerts</span>')
        if master is None:
            flags.append('<span class="flag warn">Not in master</span>')

        # Same Alert Low/High overlay as the .pptx, so the in-app gallery view is
        # verifiable at a glance too. The annotated copy lives in the repo, which
        # the /asset proxy is whitelisted to serve.
        disp_img = annotate_chart_levels(img, channel) if img and os.path.exists(img) else img
        media = (f'<img loading="lazy" src="{_asset_url(disp_img)}" alt="{_e(ticker)} chart">'
                 if img and os.path.exists(img)
                 else '<div class="noimg">chart not captured</div>')
        rows = [('Live price', fmt_num(chart.get('price')))]
        if master:
            rows += [('Holdings', fmt_pounds(master.get('holdings'))),
                     ('Alert low', fmt_num(master.get('alert_low'))),
                     ('Alert high', fmt_num(master.get('alert_high')))]
        if channel and channel.get('kind') not in (None, 'rejected'):
            rows.append(('Channel', f"{channel['kind']} {fmt_num(channel.get('lower'))}–{fmt_num(channel.get('upper'))}"))
        rows.append(('TV alerts', str(len(tv_alerts))))
        stat_html = ''.join(
            f'<div class="kv"><span>{_e(k)}</span><b>{_e(v)}</b></div>' for k, v in rows)
        tv_link = (f'<a class="tvlink" href="{TV_LAYOUT_URL.format(chart_id=cid)}" target="_blank" '
                   f'rel="noopener">Open TradingView layout ↗</a>' if cid else '')
        return (f'<article class="card">{media}'
                f'<div class="cardbody"><header><h3>{_e(ticker)}</h3>'
                f'<span class="co">{_e(chart.get("description") or "")}</span></header>'
                f'<div class="flags">{"".join(flags)}</div>'
                f'<div class="kvs">{stat_html}</div>{tv_link}</div></article>')

    sections = []
    current = None
    buf = []
    for chart in charts:
        if chart['name'] != current:
            if buf:
                sections.append(f'<section class="layout"><h2>{_e(current)}</h2>'
                                f'<div class="grid">{"".join(buf)}</div></section>')
                buf = []
            current = chart['name']
        buf.append(card(chart))
    if buf:
        sections.append(f'<section class="layout"><h2>{_e(current)}</h2>'
                        f'<div class="grid">{"".join(buf)}</div></section>')

    chips = [
        ('Layouts', stats['layouts'], 'ok'),
        ('Charts', stats['charts'], 'ok'),
        ('Missing charts', stats['missing_charts'], 'bad' if stats['missing_charts'] else 'ok'),
        ('No alerts', stats['no_alerts'], 'bad' if stats['no_alerts'] else 'ok'),
        ('Held, no chart', stats['held_no_chart'], 'bad' if stats['held_no_chart'] else 'ok'),
    ]
    chip_html = ''.join(
        f'<div class="chip {tone}"><b>{v}</b><span>{_e(label)}</span></div>' for label, v, tone in chips)

    doc = f"""<!DOCTYPE html><html lang="en"><head><meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>Investment Review Deck</title><style>
:root{{--bg:#0e1526;--panel:#16203a;--line:#263353;--text:#e8ecf5;--muted:#8a97b4;
--accent:#e0a03d;--good:#34b27b;--bad:#e5544b;--warn:#e0a03d;}}
*{{box-sizing:border-box}}body{{margin:0;background:var(--bg);color:var(--text);
font-family:system-ui,Segoe UI,Arial,sans-serif;line-height:1.45}}
header.top{{padding:28px 32px 8px}}h1{{margin:0;font-size:24px;letter-spacing:.01em}}
.sub{{color:var(--muted);font-size:13px;margin-top:4px}}
.chips{{display:flex;flex-wrap:wrap;gap:10px;padding:16px 32px 8px}}
.chip{{background:var(--panel);border:1px solid var(--line);border-radius:10px;
padding:10px 14px;min-width:96px;display:flex;flex-direction:column;gap:2px}}
.chip b{{font-size:22px;font-variant-numeric:tabular-nums}}
.chip span{{font-size:11px;text-transform:uppercase;letter-spacing:.06em;color:var(--muted)}}
.chip.bad b{{color:var(--bad)}}.chip.ok b{{color:var(--good)}}
main{{padding:12px 32px 60px}}
.layout{{margin-top:28px}}.layout h2{{font-size:13px;text-transform:uppercase;
letter-spacing:.08em;color:var(--accent);border-bottom:1px solid var(--line);
padding-bottom:8px;margin:0 0 16px}}
.grid{{display:grid;grid-template-columns:repeat(auto-fill,minmax(300px,1fr));gap:16px}}
.card{{background:var(--panel);border:1px solid var(--line);border-radius:12px;
overflow:hidden;display:flex;flex-direction:column}}
.card img{{width:100%;aspect-ratio:16/10;object-fit:cover;background:#0b0f18;display:block}}
.noimg{{aspect-ratio:16/10;display:flex;align-items:center;justify-content:center;
color:var(--muted);background:repeating-linear-gradient(45deg,#141d33,#141d33 10px,#111a2e 10px,#111a2e 20px);font-size:13px}}
.cardbody{{padding:14px 16px 16px;display:flex;flex-direction:column;gap:10px}}
.card header{{display:flex;align-items:baseline;gap:8px;flex-wrap:wrap}}
.card h3{{margin:0;font-size:17px}}.co{{color:var(--muted);font-size:12px}}
.flags{{display:flex;gap:6px;flex-wrap:wrap}}.flags:empty{{display:none}}
.flag{{font-size:10px;font-weight:700;text-transform:uppercase;letter-spacing:.05em;
padding:3px 7px;border-radius:5px}}
.flag.bad{{background:rgba(229,84,75,.16);color:#ff8a82;border:1px solid rgba(229,84,75,.4)}}
.flag.warn{{background:rgba(224,160,61,.16);color:var(--accent);border:1px solid rgba(224,160,61,.4)}}
.kvs{{display:grid;grid-template-columns:1fr 1fr;gap:6px 14px}}
.kv{{display:flex;justify-content:space-between;font-size:12px;border-bottom:1px dotted var(--line);padding-bottom:3px}}
.kv span{{color:var(--muted)}}.kv b{{font-variant-numeric:tabular-nums}}
.tvlink{{color:var(--accent);font-size:12px;text-decoration:none;font-weight:600}}
.tvlink:hover{{text-decoration:underline}}
@media(prefers-color-scheme:light){{:root{{--bg:#f5f7fb;--panel:#fff;--line:#dde3ee;
--text:#16203a;--muted:#5b6684}}}}
</style></head><body>
<header class="top"><h1>Investment Review Deck</h1>
<div class="sub">{_e(stats['charts'])} charts across {_e(stats['layouts'])} layouts · built {_e(datetime.now().strftime('%Y-%m-%d %H:%M'))}</div></header>
<div class="chips">{chip_html}</div>
<main>{''.join(sections)}</main></body></html>"""

    os.makedirs(PIPELINE_APP_DIR, exist_ok=True)
    with open(path, 'w', encoding='utf-8') as f:
        f.write(doc)


def write_summary_json(path, stats, out_pptx):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with open(path, 'w', encoding='utf-8') as f:
        json.dump({'built_at': datetime.now().isoformat(timespec='seconds'),
                   'pptx': out_pptx, **stats}, f, indent=2)


def main():
    out_path = sys.argv[1] if len(sys.argv) > 1 else os.path.join(DOWNLOADS, 'Investment_Review_Deck.pptx')

    charts = load_json('layout_manifest_tmp.json') or []
    alerts = load_json('alerts_manifest_tmp.json') or []
    channels = {c['ticker']: c for c in (load_json('channel_results_tmp.json') or []) if c.get('ticker')}
    analyst_notes = load_json('analyst_notes.json') or {}
    master_index = load_master_index()

    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H

    # ── summary page ────────────────────────────────────────────────────────
    missing_charts = [c['ticker'] for c in charts
                      if not (c.get('screenshot') and os.path.exists(c['screenshot']))]
    no_alert_tickers = []
    seen = set()
    for c in charts:
        t = c.get('ticker')
        if not t or t in seen:
            continue
        seen.add(t)
        _, master = find_master(master_index, t)
        if not alerts_for(alerts, t) and not (master and isinstance(master.get('alert_low'), (int, float))):
            no_alert_tickers.append(t)
    charted_master = set()
    for c in charts:
        key, m = find_master(master_index, c.get('ticker') or '')
        if m:
            charted_master.add(key)
    unchartered = sorted(k for k in master_index if k not in charted_master)
    # A missing chart matters most where money is actually held; the rest of the
    # master sheet is watchlist coverage.
    unchartered_held = [k for k in unchartered
                        if isinstance(master_index[k].get('holdings'), (int, float))
                        and master_index[k]['holdings'] > 0]
    unchartered_watch = [k for k in unchartered if k not in set(unchartered_held)]

    layout_names = []
    for c in charts:
        if c['name'] not in layout_names:
            layout_names.append(c['name'])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide, Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.7),
             [(f'Investment Review Deck — {datetime.now():%Y-%m-%d %H:%M}', 28, True, NAVY)])
    summary = [
        (f'{len(layout_names)} layouts, {len(charts)} charts, {len(alerts)} TradingView alerts, '
         f'{len(master_index)} master-sheet rows', 16, False, GREY),
        ('', 12, False, GREY),
        (f'Charts missing a screenshot: {len(missing_charts)}'
         + (f"  ({', '.join(missing_charts)})" if missing_charts else ' ✓'),
         16, bool(missing_charts), RED if missing_charts else GREEN),
        (f'Charted tickers with NO alert anywhere: {len(no_alert_tickers)}'
         + (f"  ({', '.join(no_alert_tickers)})" if no_alert_tickers else ' ✓'),
         16, bool(no_alert_tickers), RED if no_alert_tickers else GREEN),
        (f'HELD investments (holdings > £0) with NO chart: {len(unchartered_held)}'
         + (f"  ({', '.join(unchartered_held)})" if unchartered_held else ' ✓'),
         16, bool(unchartered_held), RED if unchartered_held else GREEN),
        (f'Watchlist-only master rows with no chart: {len(unchartered_watch)} '
         '(full list on the appendix slide at the end)', 14, False, GREY),
    ]
    add_text(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5), summary)

    # ── one slide per chart, grouped by layout in capture order ────────────
    current_layout = None
    for chart in charts:
        if chart['name'] != current_layout:
            current_layout = chart['name']
            n = sum(1 for c in charts if c['name'] == current_layout)
            section_slide(prs, current_layout, f'{n} chart(s)')
        _, master = find_master(master_index, chart.get('ticker') or '')
        t = (chart.get('ticker') or '').upper()
        chart_slide(prs, chart, current_layout, None, master,
                    channels.get(chart.get('ticker')), alerts_for(alerts, chart.get('ticker') or ''),
                    analyst=analyst_notes.get(t))

    # ── appendix: master rows without a chart ───────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6),
             [('Appendix — master-sheet rows with no TradingView chart', 22, True, NAVY)])
    appendix = []
    if unchartered_held:
        appendix.append(('Held investments (needs a chart adding):', 15, True, RED))
        for k in unchartered_held:
            m = master_index[k]
            appendix.append((f"  {k} — {m.get('share_name') or ''} "
                             f"(holdings {fmt_pounds(m.get('holdings'))})", 13, False, GREY))
        appendix.append(('', 10, False, GREY))
    appendix.append((f'Watchlist-only ({len(unchartered_watch)}):', 15, True, NAVY))
    appendix.append((', '.join(unchartered_watch) or '—', 10, False, GREY))
    add_text(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(6.1), appendix)

    purged = purge_old_versions(out_path)
    if purged:
        print(f'Recycled {len(purged)} old "(N)" copy(ies) of {os.path.basename(out_path)}')
    prs.save(out_path)
    n_slides = len(prs.slides._sldIdLst)
    print(f'Deck: {n_slides} slides -> {out_path}')
    print(f'Missing charts: {len(missing_charts)}; no-alert tickers: {len(no_alert_tickers)}; '
          f'held-with-no-chart: {len(unchartered_held)}; watchlist-no-chart: {len(unchartered_watch)}')

    # In-Chrome gallery + machine-readable summary for the Investment Production
    # Centre front end (served at /deck and consumed by the output bay).
    stats = {
        'layouts': len(layout_names), 'charts': len(charts), 'alerts': len(alerts),
        'master_rows': len(master_index), 'missing_charts': len(missing_charts),
        'no_alerts': len(no_alert_tickers), 'held_no_chart': len(unchartered_held),
        'watch_no_chart': len(unchartered_watch),
        'missing_chart_tickers': missing_charts, 'no_alert_tickers': no_alert_tickers,
        'held_no_chart_tickers': unchartered_held,
    }
    try:
        write_gallery(GALLERY_PATH, charts, alerts, channels, master_index, stats)
        write_summary_json(SUMMARY_PATH, stats, out_path)
        print(f'Gallery: {GALLERY_PATH}')
    except Exception as e:
        print(f'(gallery/summary not written: {e})')


if __name__ == '__main__':
    main()

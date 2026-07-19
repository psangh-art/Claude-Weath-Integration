#!/usr/bin/env python3
"""One-off (2026-07-13, user-requested): add iPhone and iPad as access devices
for the Finance Google Sheet on slide 1 of
Financial_Data_Pipeline_Architecture.pptx — two nodes above the existing
'Google Sheets' node, each with a connector line down into it, styled to match
the diagram's existing nodes (rounded rectangle, 3D4A6B fill, white Calibri
title + CADSFC-blue subtitle). Skips cleanly if the nodes are already there.

Usage: python add_mobile_access_2026-07-13.py [deck.pptx]
"""
import shutil
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

if sys.platform == "win32":
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

DECK = r"C:\Users\Paul\Downloads\Financial_Data_Pipeline_Architecture.pptx"

NODE_FILL = RGBColor(0x3D, 0x4A, 0x6B)   # same as the existing diagram nodes
TITLE_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
SUB_COLOR = RGBColor(0xCA, 0xDC, 0xFC)

# The 'Google Sheets' node sits at (10.67, 2.24) 1.55x0.50; the space above it
# (right column, y < 2.2) is empty. Two device nodes side by side, centred
# over it, with vertical connectors dropping into its top edge.
DEVICES = [
    ('iPhone', 'Views the Finance sheet', 10.17),
    ('iPad', 'Views the Finance sheet', 11.52),
]
NODE_W, NODE_H, NODE_TOP = 1.20, 0.50, 1.30
GS_TOP = 2.24


def add_device_node(slide, name, sub, left):
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(NODE_TOP),
        Inches(NODE_W), Inches(NODE_H))
    box.fill.solid(); box.fill.fore_color.rgb = NODE_FILL
    box.line.color.rgb = NODE_FILL; box.line.width = Emu(12700)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = tf.margin_bottom = Emu(18000)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = name
    r.font.name = 'Calibri'; r.font.size = Pt(8); r.font.bold = True
    r.font.color.rgb = TITLE_COLOR
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = sub
    r2.font.name = 'Calibri'; r2.font.size = Pt(6.5)
    r2.font.color.rgb = SUB_COLOR

    mid_x = Inches(left + NODE_W / 2)
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, mid_x, Inches(NODE_TOP + NODE_H), mid_x, Inches(GS_TOP))
    line.line.color.rgb = NODE_FILL
    line.line.width = Emu(15875)  # matches the diagram's other connectors


def main():
    deck_path = sys.argv[1] if len(sys.argv) > 1 else DECK
    prs = Presentation(deck_path)
    slide = prs.slides[0]

    existing = {sh.text_frame.text.split('\n')[0]
                for sh in slide.shapes if sh.has_text_frame}
    if any(name in existing for name, _, _ in DEVICES):
        print('iPhone/iPad nodes already present on slide 1 — nothing to do.')
        return

    backup = deck_path + '.bak-before-mobile-access-2026-07-13'
    shutil.copyfile(deck_path, backup)
    print(f'Backup -> {backup}')

    for name, sub, left in DEVICES:
        add_device_node(slide, name, sub, left)
        print(f'Added {name} node')

    prs.save(deck_path)
    print(f'Saved -> {deck_path}')


if __name__ == '__main__':
    main()

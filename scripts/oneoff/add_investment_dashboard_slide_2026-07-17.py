#!/usr/bin/env python3
"""Reflect the Investment Dashboard into the architecture deck
(Financial_Data_Pipeline_Architecture.pptx), user request 2026-07-17.

Adds ONE dedicated, readable slide (inserted right after the main data-flow
diagram) describing the Investment Dashboard: what it is, what it reads, its six
screens, and a big click-through button to http://localhost:4600. Click-through
hyperlinks are added on the Open button and on each screen card. Deliberately a
CLEAN standalone slide rather than a node crammed into the 131-shape flow diagram
— the text is a legible 10-14pt so the deck stays readable (the user asked for
readability), and there's no risk of overlapping the hand-tuned diagram.

Idempotent: if the slide is already present it does nothing. Backs the deck up
first, and recycles any stale 'X (N).pptx' duplicate copies (purge_old_versions).

Usage: python add_investment_dashboard_slide_2026-07-17.py [deck.pptx]
"""
import os
import shutil
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from config import downloads_file, purge_old_versions

if sys.platform == "win32":
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

DECK = downloads_file('architecturePptx')
MARKER = 'Investment Dashboard'          # first-line title used for the idempotency check

# Palette — matches the existing diagram nodes.
NAVY = RGBColor(0x2E, 0x50, 0x77)
NODE_FILL = RGBColor(0x3D, 0x4A, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUB_BLUE = RGBColor(0xCA, 0xDC, 0xFC)
INK = RGBColor(0x33, 0x33, 0x33)
ACCENT = RGBColor(0x1F, 0x7A, 0x33)       # green "open" button

DASH_URL = 'http://localhost:4600'

SCREENS = [
    ('Overview',     'Portfolio value, gains, dividends, cash, alert status'),
    ('Portfolio',    'Every holding + income funds, alert proximity, P&L'),
    ('Historic',     'Realised sells, 2026 trading profit, win rate'),
    ('Watchlist',    "'At lower boundary' buy zone · 1-day change · live-price refresh"),
    ('Intelligence', 'FTSE / DAX / STOXX50 / S&P / Nasdaq / Dow — live'),
    ('Design',       'Reorder every widget on every screen'),
]


def _set_text(tf, runs, anchor=MSO_ANCHOR.TOP):
    """runs: list of (text, size_pt, bold, color, align) paragraphs."""
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(45000)
    tf.margin_top = tf.margin_bottom = Emu(30000)
    for i, (text, size, bold, color, align) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = text
        r.font.name = 'Calibri'; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color


def add_card(slide, left, top, w, h, title, desc, url=None):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(left), Inches(top), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = NODE_FILL
    box.line.color.rgb = NODE_FILL; box.line.width = Emu(12700)
    _set_text(box.text_frame, [
        (title, 12, True, WHITE, PP_ALIGN.LEFT),
        (desc, 9, False, SUB_BLUE, PP_ALIGN.LEFT),
    ], anchor=MSO_ANCHOR.MIDDLE)
    if url:
        box.click_action.hyperlink.address = url
    return box


def build_slide(prs):
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank)

    # Title + subtitle
    t = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.5))
    _set_text(t.text_frame, [('Investment Dashboard', 26, True, NAVY, PP_ALIGN.LEFT)])
    s = slide.shapes.add_textbox(Inches(0.4), Inches(0.72), Inches(12.5), Inches(0.5))
    _set_text(s.text_frame, [(
        'Local web app (scripts/dashboard_app + dashboard_server.js, port 4600) — reads the '
        'master workbook and history.db, and presents the portfolio, watchlist and live market '
        'intelligence. Auto-refreshed by every pipeline run (regenerated + live-notified as the '
        'final stage of the flow), so it always reflects the latest data.',
        12, False, INK, PP_ALIGN.LEFT)])

    # Big click-through "open" button
    btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.4), Inches(1.5), Inches(4.6), Inches(0.62))
    btn.fill.solid(); btn.fill.fore_color.rgb = ACCENT
    btn.line.color.rgb = ACCENT
    _set_text(btn.text_frame, [('▶  Open dashboard  —  ' + DASH_URL, 13, True, WHITE, PP_ALIGN.CENTER)],
              anchor=MSO_ANCHOR.MIDDLE)
    btn.click_action.hyperlink.address = DASH_URL

    # Data-sources note (to the right of the button)
    d = slide.shapes.add_textbox(Inches(5.3), Inches(1.5), Inches(7.6), Inches(0.62))
    _set_text(d.text_frame, [
        ('Data sources', 11, True, NAVY, PP_ALIGN.LEFT),
        ('Stocks_Buy_Strategy.xlsx (Investments, Stocks of Interest) · history.db (captured prices) · '
         'Yahoo Finance chart API for the live indices + live watchlist prices', 9.5, False, INK, PP_ALIGN.LEFT),
    ])

    # Section label
    lbl = slide.shapes.add_textbox(Inches(0.4), Inches(2.45), Inches(12.5), Inches(0.3))
    _set_text(lbl.text_frame, [('Screens', 13, True, NAVY, PP_ALIGN.LEFT)])

    # Screen cards — 3 across, 2 rows, generously sized for legibility.
    card_w, card_h, gap = 4.05, 1.55, 0.28
    x0, y0 = 0.4, 2.85
    for i, (name, desc) in enumerate(SCREENS):
        col, row = i % 3, i // 3
        left = x0 + col * (card_w + gap)
        top = y0 + row * (card_h + gap)
        add_card(slide, left, top, card_w, card_h, name, desc, url=DASH_URL)

    # Footer note
    f = slide.shapes.add_textbox(Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.4))
    _set_text(f.text_frame, [(
        'Every card and the Open button link to the running dashboard at ' + DASH_URL +
        '. Refresh (top-right) rebuilds from the latest captured data; the Intelligence ↻ pulls live index quotes.',
        10, False, INK, PP_ALIGN.LEFT)])
    return slide


def move_after_first(prs):
    """Put the just-appended slide immediately after the main diagram (index 1)."""
    sld_lst = prs.slides._sldIdLst
    ids = list(sld_lst)
    sld_lst.remove(ids[-1])
    sld_lst.insert(1, ids[-1])


def main():
    deck_path = sys.argv[1] if len(sys.argv) > 1 else DECK
    prs = Presentation(deck_path)

    backup = deck_path + '.bak-before-dashboard-slide-2026-07-17'
    shutil.copyfile(deck_path, backup)
    print(f'Backup -> {backup}')

    # Re-runnable REPLACE (was a skip-if-present no-op). Two ordering rules matter:
    #  1) ADD the new slide FIRST, THEN delete the old one. add_slide() derives the new
    #     part name from the current slide COUNT — deleting first frees a gap that makes
    #     the new slide collide with an existing part (the agents slide became
    #     "slide4.xml" twice -> "Duplicate name" on save, and the agents slide was lost).
    #  2) When deleting, drop the RELATIONSHIP as well as the sldId, or the orphaned
    #     part's stale rels resurrect it on save.
    build_slide(prs)                       # appended at the end with a fresh part name
    move_after_first(prs)                  # move the new slide to index 1

    sld_lst = prs.slides._sldIdLst
    sld_ids = list(sld_lst)
    slides = list(prs.slides)
    for i, s in enumerate(slides):
        if i == 1:
            continue                       # skip the just-added slide (now at index 1)
        if any(sh.has_text_frame and sh.text_frame.text.split('\n')[0].strip() == MARKER
               for sh in s.shapes):
            sid = sld_ids[i]
            prs.part.drop_rel(sid.rId)
            sld_lst.remove(sid)
            print('Removed the previous Investment Dashboard slide (replaced).')
            break

    purged = purge_old_versions(deck_path)
    if purged:
        print(f'Recycled {len(purged)} old "(N)" copy(ies) of {os.path.basename(deck_path)}')
    prs.save(deck_path)
    print(f'Saved -> {deck_path} ({len(prs.slides._sldIdLst)} slides)')


if __name__ == '__main__':
    main()

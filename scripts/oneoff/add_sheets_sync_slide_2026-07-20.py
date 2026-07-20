#!/usr/bin/env python3
"""Add the Finance Google Sheet sync to the architecture deck
(Financial_Data_Pipeline_Architecture.pptx), user request 2026-07-20.

ONE readable slide covering how the master workbook reaches the user's Finance
Google Sheet: the two routes (today's manual browser import vs the service-account
API sync being built), the one-time credential setup, and the two environment
gotchas that cost a session to diagnose — the background-tab file-picker block and
Norton's TLS interception.

Same conventions as the other deck editors: legible 9-14pt on a clean standalone
slide rather than a node crammed into the 131-shape flow diagram, backs the deck up
first, recycles stale 'X (N).pptx' copies, and is RE-RUNNABLE (adds the new slide,
then removes the previous copy of it — in that order; see the ordering note in
add_investment_dashboard_slide_2026-07-17.py, deleting first collides part names).

Usage: python scripts/oneoff/add_sheets_sync_slide_2026-07-20.py [deck.pptx]
"""
import os
import shutil
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), '..'))  # scripts/
from config import downloads_file, purge_old_versions, CFG

if sys.platform == "win32":
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

DECK = downloads_file('architecturePptx')
MARKER = 'Finance Google Sheet sync'      # first-line title used for the replace check

NAVY = RGBColor(0x2E, 0x50, 0x77)
NODE_FILL = RGBColor(0x3D, 0x4A, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUB_BLUE = RGBColor(0xCA, 0xDC, 0xFC)
INK = RGBColor(0x33, 0x33, 0x33)
WARN = RGBColor(0x8A, 0x2B, 0x2B)
WARN_FILL = RGBColor(0xFB, 0xEE, 0xEE)
OK_FILL = RGBColor(0xEC, 0xF6, 0xEE)
OK_EDGE = RGBColor(0x1F, 0x7A, 0x33)

SHEET_URL = f"https://docs.google.com/spreadsheets/d/{CFG['financeSheetId']}/edit"

# The one-time credential setup, in the order it must be done. Step 5 is the one
# people skip: a valid key still 403s until the sheet is shared with the account.
SETUP = [
    ('1 · Cloud project',
     'console.cloud.google.com — project finance-sheet-sync'),
    ('2 · Enable Sheets API',
     'Sheets API only; Drive API is not needed to edit an existing sheet'),
    ('3 · Service account',
     'finance-sync@… — no project roles needed, they grant Cloud access, not sheet access'),
    ('4 · JSON key',
     r'saved to C:\Users\Paul\.secrets\ — OUTSIDE the repo, which syncs to OneDrive'),
    ('5 · Share the sheet',
     'add the service-account email as Editor — WITHOUT this every call 403s'),
    ('6 · Verify',
     'scripts/check_sheets_auth.py — read-only, names which of the three failed'),
]


def _set_text(tf, runs, anchor=MSO_ANCHOR.TOP):
    """runs: list of (text, size_pt, bold, color, align) paragraphs."""
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(45000)
    tf.margin_top = tf.margin_bottom = Emu(30000)
    for i, (text, size, bold, color, align) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = text
        r.font.name = 'Calibri'; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color


def add_card(slide, left, top, w, h, title, desc):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(left), Inches(top), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = NODE_FILL
    box.line.color.rgb = NODE_FILL; box.line.width = Emu(12700)
    _set_text(box.text_frame, [
        (title, 11.5, True, WHITE, PP_ALIGN.LEFT),
        (desc, 8.5, False, SUB_BLUE, PP_ALIGN.LEFT),
    ], anchor=MSO_ANCHOR.MIDDLE)
    return box


def add_panel(slide, left, top, w, h, fill, edge, runs):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(left), Inches(top), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.color.rgb = edge; box.line.width = Emu(12700)
    _set_text(box.text_frame, runs)
    return box


def build_slide(prs):
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank)

    t = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.5))
    _set_text(t.text_frame, [(MARKER, 26, True, NAVY, PP_ALIGN.LEFT)])

    s = slide.shapes.add_textbox(Inches(0.4), Inches(0.72), Inches(12.5), Inches(0.6))
    _set_text(s.text_frame, [(
        'Stocks_Buy_Strategy.xlsx is always the source of truth; the Google-side data tabs are '
        'derived and must never be hand-edited. The sheet is what the user reads on his phone, so '
        'it has to match the workbook after every pipeline run.',
        12, False, INK, PP_ALIGN.LEFT)])

    # ── The two routes, side by side ──────────────────────────────────────────
    add_panel(slide, 0.4, 1.45, 6.15, 1.62, WARN_FILL, WARN, [
        ('Route in use today — manual browser import', 12, True, WARN, PP_ALIGN.LEFT),
        ('Delete all data tabs (keeping ClaudeCode), then File → Import → Upload → '
         'Insert new sheet(s) + Import theme. Carries the xlsx formatting across.', 9, False, INK, PP_ALIGN.LEFT),
        ('HAZARD: tabs are deleted BEFORE the import, so a failure leaves the live sheet EMPTY. '
         'Cannot be automated — see the file-picker note below.', 9, True, WARN, PP_ALIGN.LEFT),
    ])
    add_panel(slide, 6.75, 1.45, 6.15, 1.62, OK_FILL, OK_EDGE, [
        ('Route being built — Sheets API, service account', 12, True, OK_EDGE, PP_ALIGN.LEFT),
        ('gspread + google.oauth2 write each tab IN PLACE via values.batchUpdate. No deletion, '
         'no file picker, no window where the sheet is empty.', 9, False, INK, PP_ALIGN.LEFT),
        ('Runs unattended, so it can become a quiet sub-step at the end of run_full_pipeline.js. '
         'Open question: whether it also pushes formatting, or the tabs are styled once.', 9, False, INK, PP_ALIGN.LEFT),
    ])

    lbl = slide.shapes.add_textbox(Inches(0.4), Inches(3.2), Inches(12.5), Inches(0.3))
    _set_text(lbl.text_frame, [
        ('One-time credential setup  ·  a service account, not OAuth — an OAuth app in Testing mode '
         'expires its refresh token every 7 days and would break the sync weekly',
         12, True, NAVY, PP_ALIGN.LEFT)])

    # Vertical budget on a 7.5in slide is tight: two card rows + the traps panel +
    # the footer must all fit. Cards end at y0 + 2*card_h + gap = 5.94, leaving the
    # panel 6.05-7.05 and the footer 7.10-7.38. Change card_h and the panel moves too.
    card_w, card_h, gap = 4.05, 1.03, 0.28
    x0, y0 = 0.4, 3.6
    for i, (name, desc) in enumerate(SETUP):
        col, row = i % 3, i // 3
        add_card(slide, x0 + col * (card_w + gap), y0 + row * (card_h + gap),
                 card_w, card_h, name, desc)

    # ── The two environment gotchas ───────────────────────────────────────────
    add_panel(slide, 0.4, 6.05, 12.5, 1.00, WARN_FILL, WARN, [
        ('Two environment traps, both diagnosed 2026-07-20 — neither is a Google problem',
         11, True, WARN, PP_ALIGN.LEFT),
        ('FILE PICKER: Chrome will not open a native file picker for a tab whose visibilityState is '
         '"hidden", and the Claude browser extension drives tabs in the background — so the manual '
         'import cannot be automated, and the tab must be foregrounded by hand. Check visibilityState '
         'BEFORE deleting any tab.', 8.5, False, INK, PP_ALIGN.LEFT),
        ('TLS: Norton Web/Mail Shield re-signs every certificate with its own root, which it installs '
         'in the WINDOWS store — so Chrome and Node are fine but Python (verifying against certifi) '
         'fails CERTIFICATE_VERIFY_FAILED. scripts/ssl_certs.py injects truststore so Python uses the '
         'OS store. Never verify=False.', 8.5, False, INK, PP_ALIGN.LEFT),
    ])

    f = slide.shapes.add_textbox(Inches(0.4), Inches(7.10), Inches(12.5), Inches(0.28))
    _set_text(f.text_frame, [(
        'Finance sheet: ' + SHEET_URL + '   ·   permanent ClaudeCode tab is never deleted — '
        'it is the run log and satisfies Google’s "at least one sheet" rule.',
        9, False, INK, PP_ALIGN.LEFT)])
    f.click_action.hyperlink.address = SHEET_URL
    return slide


def move_to(prs, index):
    sld_lst = prs.slides._sldIdLst
    ids = list(sld_lst)
    sld_lst.remove(ids[-1])
    sld_lst.insert(index, ids[-1])


def main():
    deck_path = sys.argv[1] if len(sys.argv) > 1 else DECK
    prs = Presentation(deck_path)

    backup = deck_path + '.bak-before-sheets-sync-slide-2026-07-20'
    shutil.copyfile(deck_path, backup)
    print(f'Backup -> {backup}')

    # ADD first, THEN delete the old copy — add_slide() derives the part name from
    # the slide COUNT, so deleting first frees a gap and the new slide collides with
    # an existing part ("Duplicate name" on save).
    new_index = min(2, len(prs.slides._sldIdLst))   # after the flow diagram + dashboard slide
    build_slide(prs)
    move_to(prs, new_index)

    sld_lst = prs.slides._sldIdLst
    sld_ids = list(sld_lst)
    for i, s in enumerate(list(prs.slides)):
        if i == new_index:
            continue
        if any(sh.has_text_frame and sh.text_frame.text.split('\n')[0].strip() == MARKER
               for sh in s.shapes):
            sid = sld_ids[i]
            prs.part.drop_rel(sid.rId)      # drop the REL too, or stale rels resurrect it
            sld_lst.remove(sid)
            print('Removed the previous Finance Google Sheet sync slide (replaced).')
            break

    purged = purge_old_versions(deck_path)
    if purged:
        print(f'Recycled {len(purged)} old "(N)" copy(ies) of {os.path.basename(deck_path)}')
    prs.save(deck_path)
    print(f'Saved -> {deck_path} ({len(prs.slides._sldIdLst)} slides)')


if __name__ == '__main__':
    main()

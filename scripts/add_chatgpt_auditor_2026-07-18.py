#!/usr/bin/env python3
"""One-off (2026-07-18, user-requested): add a ChatGPT node to slide 1 of
Financial_Data_Pipeline_Architecture.pptx, labelled "System & Code auditor" — an
external independent reviewer of the repo's architecture and scripts (it already
appears as agent #8 on the Agents slide; this puts it on the flow diagram too).

Placed in the empty space to the right of the Claude Code / GitHub repo nodes,
with a connector to the GitHub repo (the code it audits). Distinct green fill so
it reads as an external auditor rather than a data source or a Claude-system node.
Skips cleanly if it's already there. Recycles stale "(N)" duplicate copies first.

Usage: python add_chatgpt_auditor_2026-07-18.py [deck.pptx]
"""
import os
import shutil
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from config import purge_old_versions

if sys.platform == "win32":
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

DECK = r"C:\Users\Paul\Downloads\Financial_Data_Pipeline_Architecture.pptx"

AUDIT_FILL = RGBColor(0x1F, 0x7A, 0x33)   # deck green = external auditor / review
TITLE_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
SUB_COLOR = RGBColor(0xD6, 0xF0, 0xDD)

# Empty area to the right of GitHub repo (verified clear): L 4.0-5.6, T 1.0-1.9.
NODE_LEFT, NODE_TOP, NODE_W, NODE_H = 4.20, 1.05, 1.55, 0.62
# GitHub repo node right-edge / centre (its coloured box sits at ~L2.0..3.3, T1.57).
GH_RIGHT_X, GH_MID_Y = 3.30, 1.82


def main():
    deck_path = sys.argv[1] if len(sys.argv) > 1 else DECK
    prs = Presentation(deck_path)
    slide = prs.slides[0]

    firsts = {sh.text_frame.text.split('\n')[0]
              for sh in slide.shapes if sh.has_text_frame}
    if 'ChatGPT' in firsts:
        print('ChatGPT node already present on slide 1 — nothing to do.')
        return

    backup = deck_path + '.bak-before-chatgpt-auditor-2026-07-18'
    shutil.copyfile(deck_path, backup)
    print(f'Backup -> {backup}')

    # Connector first (behind the box): GitHub repo -> ChatGPT (audits the code).
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(GH_RIGHT_X), Inches(GH_MID_Y),
        Inches(NODE_LEFT), Inches(NODE_TOP + NODE_H / 2))
    line.line.color.rgb = AUDIT_FILL
    line.line.width = Emu(15875)

    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(NODE_LEFT), Inches(NODE_TOP),
        Inches(NODE_W), Inches(NODE_H))
    box.fill.solid(); box.fill.fore_color.rgb = AUDIT_FILL
    box.line.color.rgb = AUDIT_FILL; box.line.width = Emu(12700)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = tf.margin_bottom = Emu(18000)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = 'ChatGPT'
    r.font.name = 'Calibri'; r.font.size = Pt(9); r.font.bold = True
    r.font.color.rgb = TITLE_COLOR
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = 'System & Code auditor'
    r2.font.name = 'Calibri'; r2.font.size = Pt(6.5)
    r2.font.color.rgb = SUB_COLOR

    purged = purge_old_versions(deck_path)
    if purged:
        print(f'Recycled {len(purged)} old duplicate copy(ies) of the deck')
    prs.save(deck_path)
    print(f'Added ChatGPT "System & Code auditor" node -> {deck_path}')


if __name__ == '__main__':
    main()

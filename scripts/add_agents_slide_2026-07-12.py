#!/usr/bin/env python3
"""Adds/refreshes the 'Claude Code agents' slide in
Financial_Data_Pipeline_Architecture.pptx describing the repo's specialist
agents (.claude/agents/*.md) — who owns what, and where each one's remit stops.
Re-runnable: an existing agents slide is replaced, not duplicated — re-run this
whenever an agent is added or its remit changes. Follows slide 2's visual
language: numbered oval + rounded name pill + a scope column + a boundary
column.

Usage: python add_agents_slide_2026-07-12.py [deck.pptx]
"""
import shutil
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

if sys.platform == "win32":
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

DECK = r"C:\Users\Paul\Downloads\Financial_Data_Pipeline_Architecture.pptx"

NAVY = RGBColor(0x1F, 0x38, 0x64)
TEAL = RGBColor(0x11, 0x78, 0x64)
GREY = RGBColor(0x60, 0x60, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PILL = RGBColor(0x2E, 0x50, 0x77)

AGENTS = [
    ('app-developer',
     'Owns the Investment Production Centre front end (scripts/pipeline_app + the '
     'server’s presentation routes): stage filmstrip, feedstock panel, output-bay '
     'product links, styling, SSE event handling.',
     'Not for pipeline logic — chart capture, OCR and the Excel writes live in '
     'run_full_pipeline.js and the Python scripts.'),
    ('excel-formatter',
     'Visual formatting of the produced workbooks (Stocks_Buy_Strategy.xlsx, '
     'spending_summary.xlsx): colour-scheme consistency, background/foreground '
     'contrast, keeping new formatting on the established palette.',
     'Not for data or formula correctness — only how the workbooks look.'),
    ('test-analyst',
     'End-to-end data-quality audits: verifies numbers survive every hand-off — '
     'capture manifests → tradingview_layouts.xlsx → OCR reads → master-sheet '
     'writes → review deck → Finance Google Sheet. Read-only; reports PASS/FAIL '
     'with the exact cells that disagree.',
     'Finds and localises problems; does not fix production code.'),
    ('product-owner',
     'Owns BACKLOG.md: captures requests and review findings as items, prioritises '
     'against the user’s goals, breaks approved work into tasks and routes it to '
     'the right specialist agent (or builds it), keeps CLAUDE.md’s open items in '
     'sync.',
     'The user sets direction and merges PRs; ad-hoc one-off fixes go straight to '
     'the main assistant.'),
    ('investment-analyst',
     'BlackRock-grade stock analysis: fundamentals (cash flow, debt, valuation), '
     'buy prices from the user’s own charts (parallel-channel bottom, yellow '
     'trendline fallback), Analyst view on each deck slide, image-quality logging, '
     'and the daily market brief (oil/gold/FTSE/S&P + near-buy-point stocks) as a '
     'Gmail draft to the user.',
     'Decision support only — drafts, never sends; no trades; flags unreadable '
     'charts instead of guessing.'),
    ('data-developer',
     'The pipeline’s data layer: bank/broker CSV loaders (Amex/Barclays/Fidelity, '
     'fidelity_file_classifier, preflight_check), ticker normalisation, the spending '
     'pivots and future-month estimates, and the alert/below-alert derivations and '
     'ticker matching in update_master_sheet.py.',
     'Not for chart capture/browser automation, workbook visuals (excel-formatter), '
     'the front end, or auditing (test-analyst finds; this agent fixes).'),
    ('validation',
     'Audits the review deck’s detected patterns and Alert Low/High against the '
     'signed-off pattern rules; groups faults by root cause and fixes every instance '
     'with a batch A/B; flags charts that look like a genuinely new pattern for the '
     'user to rule on.',
     'Not for running the pipeline, workbook formatting, the front end, or '
     'fundamental analysis (investment-analyst).'),
    ('ChatGPT (external)',
     'Independent code reviewer: periodic outside review of the system’s '
     'architecture and scripts (e.g. the 2026-07-12 “Investment OS” review) — '
     'challenges design choices, flags redundancy and hard-coding, proposes '
     'next steps.',
     'Advisory only — no repo access; every recommendation is challenged '
     'against repo evidence, and only accepted items reach BACKLOG.md.'),
]

EXTERNAL = {'ChatGPT (external)'}
EXTERNAL_PILL = RGBColor(0x74, 0x8A, 0x5E)  # muted green — visually "not one of ours"


def add_text(slide, x, y, w, h, text, size, bold=False, color=GREY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Arial'
    return box


def remove_existing_agents_slide(prs):
    """Re-runnable: drop any previously-added agents slide (identified by its
    title text) so re-running replaces it instead of appending a duplicate."""
    sld_ids = prs.slides._sldIdLst
    for idx, slide in enumerate(list(prs.slides)):
        texts = [sh.text_frame.text for sh in slide.shapes if sh.has_text_frame]
        if any(t.startswith('Claude Code agents') for t in texts):
            rId = sld_ids[idx].rId
            prs.part.drop_rel(rId)
            sld_ids.remove(sld_ids[idx])
            print(f'Removed existing agents slide (was slide {idx + 1})')
            return


def main():
    deck_path = sys.argv[1] if len(sys.argv) > 1 else DECK
    backup = deck_path + '.bak-before-agents-slide-2026-07-12'
    shutil.copyfile(deck_path, backup)
    print(f'Backup -> {backup}')

    prs = Presentation(deck_path)
    remove_existing_agents_slide(prs)
    slide = prs.slides.add_slide(prs.slides[1].slide_layout)  # match slide 2's layout

    add_text(slide, Inches(0.4), Inches(0.2), Inches(9.5), Inches(0.4),
             'Claude Code agents — who builds and checks what', 18, bold=True, color=NAVY)
    add_text(slide, Inches(0.4), Inches(0.62), Inches(11.5), Inches(0.3),
             'Specialist agents in .claude/agents/ (plus ChatGPT as an outside reviewer) — '
             'each owns one part of the system; the user sets direction and merges all PRs',
             11, color=GREY)

    # Rows adapt to the agent count so they always fit above the 7.5in bottom:
    # step is capped at 1.06 (roomy for a few) and shrinks as agents are added.
    top0 = 0.98
    bottom_limit = 7.28
    step = min(1.06, (bottom_limit - top0) / max(1, len(AGENTS)))
    small = step < 0.95                 # tighten fonts/pill when rows are dense
    name_pt = 12 if small else 13
    body_pt = 9 if small else 10
    pill_h = min(0.5, step - 0.28)
    row_h = Inches(step - 0.05)
    for i, (name, scope, boundary) in enumerate(AGENTS):
        y = Inches(top0 + i * step)
        num = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.4), y + Inches(0.08), Inches(0.3), Inches(0.3))
        num.fill.solid(); num.fill.fore_color.rgb = NAVY
        num.line.fill.background()
        tf = num.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(i + 1)
        r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = 'Arial'

        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), y, Inches(2.45), Inches(pill_h))
        pill.fill.solid(); pill.fill.fore_color.rgb = EXTERNAL_PILL if name in EXTERNAL else PILL
        pill.line.fill.background()
        tf = pill.text_frame; tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = name
        r.font.size = Pt(name_pt); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = 'Consolas'

        add_text(slide, Inches(3.5), y, Inches(4.9), row_h, scope, body_pt, color=RGBColor(0, 0, 0))
        add_text(slide, Inches(8.6), y, Inches(3.6), row_h, boundary, body_pt, color=TEAL)

    add_text(slide, Inches(3.5), Inches(top0 - 0.28), Inches(4.9), Inches(0.25),
             'OWNS', 9, bold=True, color=GREY)
    add_text(slide, Inches(8.6), Inches(top0 - 0.28), Inches(3.6), Inches(0.25),
             'STOPS AT', 9, bold=True, color=GREY)

    prs.save(deck_path)
    print(f'Added agents slide ({len(AGENTS)} agents) as slide {len(prs.slides._sldIdLst)} -> {deck_path}')


if __name__ == '__main__':
    main()

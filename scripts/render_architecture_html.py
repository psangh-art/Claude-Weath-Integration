"""Render the Financial Data Pipeline Architecture .pptx as a readable, in-app
HTML page for the Investment Dashboard.

No slide renderer (LibreOffice) is available in this environment, so the raw
.pptx link only ever downloaded the file rather than opening it. Instead of a
hand-written second copy of the architecture (which would drift from the deck),
this reads the SAME .pptx the deck scripts produce and lays every slide out as
absolutely-positioned HTML boxes, using each shape's EMU geometry, fill, border
and per-run font styling. It therefore stays in sync with the deck automatically.

Layout is fully responsive with pure CSS: each slide is a `container-type:
inline-size` stage at the slide's aspect ratio, children are positioned in % of
the slide, and font sizes are expressed in `cqw` (1% of the stage width) so text
scales with the stage at any width. Lines/connectors are drawn in one SVG per
slide whose viewBox matches the slide's EMU coordinates, so they stay exact.

Called by dashboard_server.js's `/decks/architecture` route (which regenerates
when the .pptx is newer than the cached HTML). Can also be run standalone:
    python scripts/render_architecture_html.py [in.pptx] [out.html]
"""
import base64
import html
import os
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DEFAULT_PPTX = os.path.join(os.path.expanduser('~'), 'Downloads',
                            'Financial_Data_Pipeline_Architecture.pptx')
DEFAULT_OUT = os.path.join(REPO, 'scripts', 'dashboard_app', 'architecture.html')

# EMU per point (914400 EMU/inch ÷ 72 pt/inch).
EMU_PER_PT = 914400 / 72.0


def _rgb(color_format):
    """Return '#rrggbb' for a solid RGB colour, else None (theme/inherited)."""
    try:
        if color_format is not None and color_format.type is not None \
                and str(color_format.type) == 'MSO_THEME_COLOR.NOT_THEME_COLOR':
            return '#' + str(color_format.rgb)
    except Exception:
        pass
    # Fallback path: many shapes expose .rgb only when the type is an explicit RGB.
    try:
        return '#' + str(color_format.rgb)
    except Exception:
        return None


def _fill_hex(shape):
    try:
        if shape.fill.type is not None and int(shape.fill.type) == 1:  # solid
            return '#' + str(shape.fill.fore_color.rgb)
    except Exception:
        pass
    return None


def _line_hex(shape):
    try:
        ln = shape.line
        if ln.color is not None and ln.color.type is not None:
            return '#' + str(ln.color.rgb)
    except Exception:
        pass
    return None


def _flip(shape):
    """(flipH, flipV) from the shape's xfrm, defaulting to (False, False)."""
    try:
        xfrm = shape._element.spPr.find(
            '{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
        if xfrm is not None:
            return xfrm.get('flipH') == '1', xfrm.get('flipV') == '1'
    except Exception:
        pass
    return False, False


def _para_html(shape, scale_cqw):
    """Render a text frame's paragraphs as HTML lines with per-run styling.

    scale_cqw converts a point size to cqw units: pt * (100 / slide_width_pt).
    """
    if not shape.has_text_frame:
        return ''
    lines = []
    for para in shape.text_frame.paragraphs:
        align = {None: '', 1: 'left', 2: 'center', 3: 'right'}.get(
            int(para.alignment) if para.alignment is not None else None, '')
        runs_html = []
        for run in para.runs:
            txt = html.escape(run.text)
            if not txt:
                continue
            styles = []
            try:
                if run.font.size is not None:
                    styles.append(f'font-size:{run.font.size.pt * scale_cqw:.3f}cqw')
            except Exception:
                pass
            if run.font.bold:
                styles.append('font-weight:700')
            if run.font.italic:
                styles.append('font-style:italic')
            col = _rgb(run.font.color) if run.font.color is not None else None
            if col:
                styles.append(f'color:{col}')
            runs_html.append(f'<span style="{";".join(styles)}">{txt}</span>')
        if not runs_html:
            runs_html.append('&nbsp;')
        style = f'text-align:{align}' if align else ''
        lines.append(f'<div class="ln" style="{style}">{"".join(runs_html)}</div>')
    return ''.join(lines)


def _pct(v, total):
    return f'{(v / total) * 100:.4f}%'


def render(pptx_path=DEFAULT_PPTX, out_path=DEFAULT_OUT):
    prs = Presentation(pptx_path)
    sw, sh = prs.slide_width, prs.slide_height
    slide_w_pt = sw / EMU_PER_PT
    scale_cqw = 100.0 / slide_w_pt  # pt -> cqw

    slides_html = []
    for si, slide in enumerate(prs.slides):
        boxes = []
        lines = []  # svg lines in EMU coords
        for shape in slide.shapes:
            try:
                left = shape.left
                top = shape.top
                w = shape.width
                h = shape.height
            except Exception:
                continue
            if left is None or top is None or w is None or h is None:
                continue

            if shape.shape_type == MSO_SHAPE_TYPE.LINE:
                fh, fv = _flip(shape)
                x1, x2 = (left + w, left) if fh else (left, left + w)
                y1, y2 = (top + h, top) if fv else (top, top + h)
                stroke = _line_hex(shape) or '#8895a7'
                lines.append(
                    f'<line x1="{x1}" y1="{y1}" x2="{x2}" y2="{y2}" '
                    f'stroke="{stroke}" stroke-width="14000" '
                    f'stroke-linecap="round"/>')
                continue

            pos = (f'left:{_pct(left, sw)};top:{_pct(top, sh)};'
                   f'width:{_pct(w, sw)};height:{_pct(h, sh)}')

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = shape.image
                    b64 = base64.b64encode(img.blob).decode('ascii')
                    src = f'data:{img.content_type};base64,{b64}'
                    boxes.append(
                        f'<img class="pic" style="{pos}" src="{src}" alt="">')
                except Exception:
                    pass
                continue

            # AUTO_SHAPE / TEXT_BOX / anything else with a frame.
            fill = _fill_hex(shape)
            border = _line_hex(shape)
            bstyle = [pos]
            if fill:
                bstyle.append(f'background:{fill}')
            if border:
                bstyle.append(f'border:0.12cqw solid {border}')
            radius = ''
            try:
                if str(shape.shape_type) == 'ROUNDED_RECTANGLE (5)':
                    radius = 'border-radius:0.6cqw'
            except Exception:
                pass
            if radius:
                bstyle.append(radius)
            inner = _para_html(shape, scale_cqw)
            boxes.append(
                f'<div class="box" style="{";".join(bstyle)}">{inner}</div>')

        svg = ''
        if lines:
            svg = (f'<svg class="lines" viewBox="0 0 {sw} {sh}" '
                   f'preserveAspectRatio="none">{"".join(lines)}</svg>')
        slides_html.append(
            f'<section class="slidecard">'
            f'<div class="stage" style="aspect-ratio:{sw}/{sh}">'
            f'{svg}{"".join(boxes)}</div></section>')

    body = '\n'.join(slides_html)
    mtime = os.path.getmtime(pptx_path)
    page = _PAGE.replace('{{SLIDES}}', body).replace(
        '{{COUNT}}', str(len(prs.slides)))
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    with open(out_path, 'w', encoding='utf-8') as f:
        f.write(page)
    return out_path, mtime


_PAGE = """<!doctype html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>Pipeline Architecture</title>
<style>
  :root{
    --bg:#f4f6fb; --card:#ffffff; --ink:#1c2431; --muted:#5b6675;
    --line:#e2e7f0; --shadow:0 6px 22px rgba(20,30,50,.08);
  }
  @media (prefers-color-scheme: dark){
    :root{ --bg:#0e1420; --card:#151d2b; --ink:#e7ecf5; --muted:#93a0b4;
      --line:#26303f; --shadow:0 8px 26px rgba(0,0,0,.45); }
  }
  :root[data-theme="dark"]{ --bg:#0e1420; --card:#151d2b; --ink:#e7ecf5;
    --muted:#93a0b4; --line:#26303f; --shadow:0 8px 26px rgba(0,0,0,.45); }
  :root[data-theme="light"]{ --bg:#f4f6fb; --card:#ffffff; --ink:#1c2431;
    --muted:#5b6675; --line:#e2e7f0; --shadow:0 6px 22px rgba(20,30,50,.08); }
  *{ box-sizing:border-box; }
  body{ margin:0; background:var(--bg); color:var(--ink);
    font-family:'Segoe UI',system-ui,-apple-system,sans-serif;
    padding:24px clamp(12px,4vw,48px); }
  header{ max-width:1180px; margin:0 auto 18px; }
  header h1{ font-size:22px; margin:0 0 4px; }
  header p{ margin:0; color:var(--muted); font-size:13px; }
  .wrap{ max-width:1180px; margin:0 auto; display:flex; flex-direction:column;
    gap:22px; }
  .slidecard{ background:var(--card); border:1px solid var(--line);
    border-radius:14px; box-shadow:var(--shadow); padding:14px; overflow:hidden; }
  .stage{ position:relative; width:100%; container-type:inline-size;
    background:#ffffff; border-radius:8px; overflow:hidden; }
  @media (prefers-color-scheme: dark){ .stage{ background:#f7f9fc; } }
  :root[data-theme="dark"] .stage{ background:#f7f9fc; }
  .box{ position:absolute; display:flex; flex-direction:column;
    justify-content:center; padding:0.35cqw 0.55cqw; overflow:hidden;
    line-height:1.18; color:#222; }
  .box .ln{ white-space:pre-wrap; word-break:break-word; }
  .pic{ position:absolute; object-fit:contain; }
  .lines{ position:absolute; inset:0; width:100%; height:100%;
    pointer-events:none; z-index:0; }
  .box, .pic{ z-index:1; }
</style>
</head>
<body>
  <header>
    <h1>Financial Data Pipeline — Architecture</h1>
    <p>Live in-app view of the architecture deck ({{COUNT}} slides), rendered from the PowerPoint so it always matches the latest build.</p>
  </header>
  <div class="wrap">
    {{SLIDES}}
  </div>
  <script>
    // Inherit the dashboard's theme choice if the parent set one.
    try{ var t=localStorage.getItem('dash-theme');
      if(t) document.documentElement.setAttribute('data-theme',t); }catch(e){}
  </script>
</body>
</html>
"""


if __name__ == '__main__':
    src = sys.argv[1] if len(sys.argv) > 1 else DEFAULT_PPTX
    dst = sys.argv[2] if len(sys.argv) > 2 else DEFAULT_OUT
    out, _ = render(src, dst)
    print('Wrote', out)

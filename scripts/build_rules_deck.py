#!/usr/bin/env python3
"""Build the ALERT RULES deck — a visual spec of how Alert Low / Alert High are
derived from the user's TradingView drawings, one named pattern per slide with a
real captured chart as the example, followed by the charts whose pattern the rules
do NOT yet explain.

This is a review artefact, not part of the pipeline: it exists so the user can check
the agreed rules against real charts and refine them. The trailing "needs a rule"
section is the point of the deck — those are the charts to argue about.

Usage:
  python build_rules_deck.py [results.json] [out.pptx]
    results.json  channel_detect --batch output (default logs/channel_results_newrules.json)
    out.pptx      default ~/Downloads/Alert_Rules_Model.pptx
"""
import os
import sys
import json

from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, SCRIPT_DIR)

# Reuse the review deck's layout helpers and palette rather than re-inlining them —
# both decks must look like the same document.
from build_review_deck import (add_text, add_picture_fitted, section_slide,
                               NAVY, RED, AMBER, GREEN, GREY, SLIDE_W, SLIDE_H)
from config import downloads_dir, purge_old_versions
from channel_detect import ON_ALERT_TOL

DEFAULT_RESULTS = os.path.join(os.path.dirname(SCRIPT_DIR), 'logs', 'channel_results_newrules.json')
DEFAULT_OUT = os.path.join(downloads_dir(), 'Alert_Rules_Model.pptx')
INPUT_MANIFEST = os.path.join(SCRIPT_DIR, 'channel_input_tmp.json')

BLUE = RGBColor(0x1F, 0x6F, 0xD0)
YELLOW_TXT = RGBColor(0xB9, 0x77, 0x0E)

# channel_detect's pattern string for a chart with no blue channel — only hand-drawn
# yellow lines. The TREND LINES ONLY model covers these (agreed 2026-07-15).
TREND_ONLY = 'trend lines only (no blue channel read)'

# The named patterns, in the order the rules are explained. Each is keyed by the
# `pattern` string channel_detect emits, with the example chart chosen by hand as
# the clearest instance of it.
PATTERNS = [
    {
        'name': 'IN-CHANNEL',
        'match': 'price INSIDE channel',
        'example': 'SDLF',
        'shape': 'Price sits between the two rails of a drawn parallel channel.',
        'rule': [
            'Alert Low  = the rail below price (or a yellow trend line if one sits nearer).',
            'Alert High = the rail above price (or a nearer yellow line).',
        ],
    },
    {
        'name': 'BREAKOUT ABOVE',
        'match': 'price ABOVE channel (broken out) — top rail is support',
        'example': 'SDR',
        'shape': 'Price has broken out ABOVE the channel — both rails are below it.',
        'rule': [
            'The TOP rail flips role and becomes support: Alert Low = top rail.',
            'Alert High = only if a line still sits above price; otherwise left blank.',
            'The rail keeps its geometric meaning, not its original "upper" label.',
        ],
    },
    {
        'name': 'BREAKDOWN BELOW',
        'match': 'price BELOW channel (broken down) — band still governs',
        'example': 'RIO',
        'shape': 'Price has dropped out of the BOTTOM of the channel — both rails are above it.',
        'rule': [
            'The channel IS the band: Alert Low = bottom rail, Alert High = top rail.',
            'Dropping out of the bottom is the buy signal, not a reason to re-cast the rails.',
            'Alert Low sits ABOVE today\'s price, so the row flags as below Alert Low.',
            'Parallel-only: yellow lines drawn nearer to price still win.',
        ],
    },
    {
        'name': 'TREND LINES ONLY',
        'match': TREND_ONLY,
        'example': 'AZN',
        'shape': 'No blue channel — support and resistance are hand-drawn yellow lines.',
        'rule': [
            'Yellow lines are treated exactly like rails: nearest below / nearest above.',
            'A line must be dead straight to count — a wavy yellow indicator is rejected.',
        ],
    },
    {
        'name': 'WEDGE',
        'match': 'wedge (trend lines converging)',
        'example': 'HOC',
        'shape': 'Two yellow trend lines closing on each other in the near future.',
        'rule': [
            'Alert Low and Alert High are UNCHANGED — the trend-line rules still stand.',
            'The wedge is the shape, not a new way to price the alerts.',
            'A break ABOVE the wedge is a potential buy.',
            'The lines must meet within ~3 months, and still be 3% apart today.',
        ],
    },
    {
        'name': 'ON THE LINE',
        'match': 'price ON a drawn line',
        'example': 'SILVER',
        'shape': 'Price is sitting ON a drawn line — the alert condition itself.',
        'rule': [
            'Within %.2f%% of a line, the side cannot be read from the image at all.' % (ON_ALERT_TOL * 100),
            'So no side is guessed: the line becomes Alert Low and the row is flagged ON ALERT.',
            'The x1.05 buffer is skipped — price is already there, nothing to warn early about.',
        ],
    },
    {
        'name': 'NO READ',
        'match': 'no lines read',
        'example': 'UKW',
        'shape': 'No drawn lines on the chart at all.',
        'rule': [
            'Nothing is written. Silence beats a guess in a live trading sheet.',
            'This chart previously produced Alert High 129.04 — the blue BUY button,',
            'read as a trendline. Buttons and event markers are not lines.',
        ],
    },
]


def load_json(path):
    with open(path, 'r', encoding='utf-8') as f:
        return json.load(f)


def fmt(v):
    return '—' if v is None else f'{v:,.2f}'


def rule_slide(prs, pat, rec, price):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide, Inches(0.5), Inches(0.28), Inches(12.4), Inches(0.5),
             [(pat['name'], 30, True, NAVY)])
    add_text(slide, Inches(0.5), Inches(0.85), Inches(12.4), Inches(0.35),
             [(pat['shape'], 14, False, GREY)])

    if rec and rec.get('screenshot') and os.path.exists(rec['screenshot']):
        add_picture_fitted(slide, rec['screenshot'], Inches(0.5), Inches(1.35),
                           Inches(8.0), Inches(5.6))

    x = Inches(8.75)
    add_text(slide, x, Inches(1.35), Inches(4.3), Inches(0.3),
             [('THE RULE', 12, True, NAVY)])
    y = 1.75
    for line in pat['rule']:
        add_text(slide, x, Inches(y), Inches(4.4), Inches(0.6),
                 [('•  ' + line, 11, False, RGBColor(0, 0, 0))])
        y += 0.52

    if rec:
        y = max(y + 0.25, 3.7)
        add_text(slide, x, Inches(y), Inches(4.3), Inches(0.3),
                 [('THIS CHART  (%s)' % rec['ticker'], 12, True, NAVY)])
        y += 0.42
        rows = [('Live price', fmt(price), GREY)]
        if rec.get('blue_lines'):
            rows.append(('Blue rails read', ', '.join(fmt(v) for v in rec['blue_lines']), BLUE))
        if rec.get('yellow_lines'):
            rows.append(('Yellow lines read', ', '.join(fmt(v) for v in rec['yellow_lines']), YELLOW_TXT))
        if rec.get('wedge'):
            wg = rec['wedge']
            rows.append(('Wedge lines', '%s / %s' % (fmt(wg['lower_line']), fmt(wg['upper_line'])),
                         YELLOW_TXT))
            rows.append(('They meet', '%.2f pane-widths ahead of today' % wg['apex_ahead_frac'], GREY))
        rows.append(('Alert Low', '%s  (%s)' % (fmt(rec['lower']), rec.get('alert_low_src') or 'none'), GREEN))
        rows.append(('Alert High', '%s  (%s)' % (fmt(rec['upper']), rec.get('alert_high_src') or 'none'), RED))
        if rec.get('wedge') and rec['wedge']['broken_above']:
            rows.append(('Signal', 'broke ABOVE the wedge — potential buy', GREEN))
        if rec.get('on_alert'):
            rows.append(('Status', 'ON ALERT — price has reached the line', GREEN))
        for label, val, colour in rows:
            add_text(slide, x, Inches(y), Inches(1.55), Inches(0.3), [(label, 10, True, GREY)])
            add_text(slide, x + Inches(1.6), Inches(y), Inches(2.8), Inches(0.3), [(val, 10, False, colour)])
            y += 0.3


def governing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide, Inches(0.5), Inches(0.4), Inches(12.4), Inches(0.6),
             [('The governing rule', 34, True, NAVY)])
    add_text(slide, Inches(0.5), Inches(1.3), Inches(12.4), Inches(3.4), [
        ('Alert Low  =  the nearest drawn line BELOW today\'s price', 22, True, GREEN),
        ('Alert High =  the nearest drawn line ABOVE today\'s price', 22, True, RED),
        ('', 10, False, GREY),
        ('Every line competes on equal terms — blue channel rails and yellow hand-drawn', 14, False, RGBColor(0, 0, 0)),
        ('trend lines alike. Whichever is closest to price on each side wins.', 14, False, RGBColor(0, 0, 0)),
        ('', 8, False, GREY),
        ('Each line is read at TODAY\'S date (the last candle), not at the frame edge —', 14, False, RGBColor(0, 0, 0)),
        ('a sloping line is worth a different price at those two points.', 14, False, RGBColor(0, 0, 0)),
        ('', 8, False, GREY),
        ('Lines are split by WHICH SIDE OF PRICE they fall on, not by whether they were drawn', 14, False, RGBColor(0, 0, 0)),
        ('as the "upper" or "lower" rail. Break out ABOVE a channel and the top rail changes', 14, False, RGBColor(0, 0, 0)),
        ('role, becoming support.', 14, False, RGBColor(0, 0, 0)),
        ('', 8, False, GREY),
        ('One exception — a parallel-only chart that price has dropped BELOW. There the channel', 14, True, NAVY),
        ('stays the band: bottom rail to Alert Low, top rail to Alert High. See BREAKDOWN BELOW.', 14, True, NAVY),
    ])


def buffer_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide, Inches(0.5), Inches(0.4), Inches(12.4), Inches(0.6),
             [('The buffer, and its two guards', 34, True, NAVY)])
    add_text(slide, Inches(0.5), Inches(1.3), Inches(12.4), Inches(4.2), [
        ('Alert Low is written 5% ABOVE the support line, so the alert fires before', 15, False, RGBColor(0, 0, 0)),
        ('price actually reaches it.   Alert Low = support x 1.05', 15, True, NAVY),
        ('', 10, False, GREY),
        ('Guard 1 — CLAMP.  The buffer never crosses Alert High.', 16, True, GREEN),
        ('When the nearest support sits less than 5% below the nearest resistance, x1.05', 13, False, RGBColor(0, 0, 0)),
        ('walks Alert Low straight past it. That put 8 rows live with Alert Low ABOVE', 13, False, RGBColor(0, 0, 0)),
        ('Alert High (GOLD, SILVER, LAND, CPG, RIO, STJ, HIK, DCC). The buffer now yields.', 13, False, RGBColor(0, 0, 0)),
        ('', 8, False, GREY),
        ('Guard 2 — NO BUFFER ONCE REACHED, OR PASSED.', 16, True, GREEN),
        ('If price is sitting on the line, or has already gone below it, there is nothing left', 13, False, RGBColor(0, 0, 0)),
        ('to warn about — and inflating the level by 5% would put Alert Low above a price you', 13, False, RGBColor(0, 0, 0)),
        ('can buy at today. Below a parallel channel, Alert Low IS the bottom rail: RIO writes', 13, False, RGBColor(0, 0, 0)),
        ('7,054.35 against a price of 6,927 — not 7,407.07.', 13, False, RGBColor(0, 0, 0)),
        ('', 8, False, GREY),
        ('A row already holding an inverted pair is always rewritten — a broken row is not', 13, False, RGBColor(0, 0, 0)),
        ('"noise" to be left alone, which is why the old inversions survived every re-run.', 13, False, RGBColor(0, 0, 0)),
    ])


def unknown_slide(prs, rec, price, why):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide, Inches(0.5), Inches(0.28), Inches(12.4), Inches(0.5),
             [('%s — %s' % (rec['ticker'], why['headline']), 26, True, AMBER)])
    add_text(slide, Inches(0.5), Inches(0.85), Inches(12.4), Inches(0.35),
             [(why['detail'], 13, False, GREY)])

    if rec.get('screenshot') and os.path.exists(rec['screenshot']):
        add_picture_fitted(slide, rec['screenshot'], Inches(0.5), Inches(1.35),
                           Inches(8.0), Inches(5.6))

    x = Inches(8.75)
    add_text(slide, x, Inches(1.35), Inches(4.3), Inches(0.3), [('WHAT THE RULES PRODUCED', 12, True, NAVY)])
    y = 1.78
    rows = [('Live price', fmt(price), GREY),
            ('Pattern', (rec.get('pattern') or '—'), GREY)]
    if rec.get('blue_lines'):
        rows.append(('Blue rails', ', '.join(fmt(v) for v in rec['blue_lines']), BLUE))
    if rec.get('yellow_lines'):
        rows.append(('Yellow lines', ', '.join(fmt(v) for v in rec['yellow_lines']), YELLOW_TXT))
    rows += [('Alert Low', fmt(rec['lower']), GREEN), ('Alert High', fmt(rec['upper']), RED)]
    for label, val, colour in rows:
        add_text(slide, x, Inches(y), Inches(1.5), Inches(0.3), [(label, 10, True, GREY)])
        add_text(slide, x + Inches(1.55), Inches(y), Inches(2.85), Inches(0.55), [(str(val)[:70], 10, False, colour)])
        y += 0.34 if label not in ('Pattern',) else 0.62

    y += 0.2
    add_text(slide, x, Inches(y), Inches(4.3), Inches(0.3), [('THE QUESTION', 12, True, AMBER)])
    y += 0.4
    add_text(slide, x, Inches(y), Inches(4.4), Inches(1.6), [(why['question'], 11, False, RGBColor(0, 0, 0))])


def classify_unknowns(results, prices):
    """Charts the agreed rules do not convincingly explain.

    Both of the original families were retired on 2026-07-15, each because the user
    reviewed every chart it caught and confirmed the read:

    * 'implausible span' flagged any band that was a large fraction of the share
      price, on the theory that two unrelated drawings had been paired into one
      channel — PAF/IBST/BKG (trend lines only), ENT/LSEG/BREE/TW. (parallel only)
      and AUTO (a blue channel with yellow over it). A wide band is the chart, not a
      misread.
    * 'stale Alert Low' flagged RIO/CPG/STJ, where a broken-down chart yielded a bare
      resistance and left an older, higher Alert Low stranded in the sheet. Fixed at
      source instead: below a parallel channel the band still governs, so these now
      read a full pair off their rails.

    What remains is a contradiction check rather than a heuristic: a pair the rules
    produced that cannot be traded. Nothing should ever reach it — a chart appearing
    here means the rules have a hole, not that a drawing is unusual.
    """
    out = []
    for rec in results:
        lo, hi = rec.get('lower'), rec.get('upper')
        if lo is None or hi is None or lo < hi:
            continue
        out.append((rec, prices.get(rec['ticker']), {
            'headline': 'Alert Low is not below Alert High',
            'detail': 'The rules produced a pair that cannot be traded.',
            'question': ('Alert Low has come out at or above Alert High, so there is no band '
                         'between them and neither level can fire as intended. This is a gap in '
                         'the rules rather than a judgement call about the drawing — which line '
                         'was misread, and which of the two levels should stand?'),
        }))
    return out


def main():
    results_path = sys.argv[1] if len(sys.argv) > 1 else DEFAULT_RESULTS
    out_path = sys.argv[2] if len(sys.argv) > 2 else DEFAULT_OUT

    results = load_json(results_path)
    by_ticker = {r['ticker']: r for r in results}
    prices = {r['ticker']: r['known_price'] for r in load_json(INPUT_MANIFEST)}

    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide, Inches(0.8), Inches(2.4), Inches(11.7), Inches(2.0), [
        ('Alert Rules — the model', 44, True, NAVY),
        ('How Alert Low and Alert High are read from your TradingView drawings', 18, False, GREY),
        ('Named patterns, each with a real captured chart. Charts the rules do not', 14, False, GREY),
        ('yet explain are at the back — those are the ones to refine.', 14, False, GREY),
    ])

    governing_slide(prs)

    section_slide(prs, 'The patterns', 'One slide per named pattern, with a live example')
    for pat in PATTERNS:
        rec = by_ticker.get(pat['example'])
        rule_slide(prs, pat, rec, prices.get(pat['example']))

    buffer_slide(prs)

    unknowns = classify_unknowns(results, prices)
    if unknowns:
        section_slide(prs, 'Patterns we do not yet recognise',
                      '%d charts where the rules produce an answer we cannot defend' % len(unknowns))
        for rec, p, why in unknowns:
            unknown_slide(prs, rec, p, why)
    else:
        section_slide(prs, 'Every chart is explained',
                      'All %d captured reads follow the six patterns — no unresolved charts'
                      % sum(1 for r in results if r.get('kind')))

    purged = purge_old_versions(out_path)
    if purged:
        print('Recycled %d old "(N)" copy(ies) of %s' % (len(purged), os.path.basename(out_path)))
    prs.save(out_path)
    print('Rules deck -> %s' % out_path)
    print('  %d named patterns, %d charts needing a rule' % (len(PATTERNS), len(unknowns)))


if __name__ == '__main__':
    main()
